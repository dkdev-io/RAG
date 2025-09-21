#!/usr/bin/env python3
"""
Phase 2: Processing Script
Extracts content from files, performs OCR, and builds the vector index.
"""

import os
import sys
import json
import logging
import zipfile
import tarfile
from pathlib import Path
from typing import List, Dict, Any, Optional
import argparse
from datetime import datetime

import pandas as pd
import numpy as np
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from docx import Document
import faiss
from sentence_transformers import SentenceTransformer
from tqdm import tqdm
import subprocess
import shutil

import sys
sys.path.append(str(Path(__file__).parent.parent))
from config.settings import (
    DEFAULT_SOURCE_DIR, PROCESSED_DIR, CHUNK_SIZE, CHUNK_OVERLAP,
    EMBEDDING_MODEL, TESSERACT_CONFIG
)

# Use environment variables for dynamic directories
SOURCE_DIR = Path(os.getenv('RAG_SOURCE_DIR', DEFAULT_SOURCE_DIR))
ASSESSMENT_REPORT = Path(os.getenv('RAG_ASSESSMENT_REPORT', Path(__file__).parent.parent / "assessment_report.csv"))
VECTOR_DB_PATH = Path(os.getenv('RAG_VECTOR_DB_PATH', PROCESSED_DIR / "vector_db"))


# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('processing.log'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)


class FileConverter:
    """Handles conversion of incompatible file types to supported formats."""

    def __init__(self):
        self.conversion_mappings = {
            # Office formats that might need conversion
            '.pptx': '.txt',
            '.ppt': '.txt',
            '.rtf': '.txt',
            '.odt': '.txt',
            '.ods': '.txt',
            '.odp': '.txt',
            # Legacy formats
            '.wpd': '.txt',
            '.pages': '.txt',
            '.numbers': '.txt',
            '.key': '.txt',
            # Other text formats
            '.epub': '.txt',
            '.mobi': '.txt',
            # Archive formats (will extract and process contents)
            '.zip': 'extract',
            '.tar': 'extract',
            '.gz': 'extract',
            '.7z': 'extract',
            '.rar': 'extract',
        }

    def can_convert(self, file_path: Path) -> bool:
        """Check if file can be converted to a supported format."""
        extension = file_path.suffix.lower()
        return extension in self.conversion_mappings

    def convert_file(self, file_path: Path, output_dir: Path) -> List[Path]:
        """Convert file to supported format(s). Returns list of converted files."""
        extension = file_path.suffix.lower()
        target_format = self.conversion_mappings.get(extension)

        if not target_format:
            return []

        try:
            if target_format == 'extract':
                return self._extract_archive(file_path, output_dir)
            elif target_format == '.txt':
                return self._convert_to_text(file_path, output_dir)
            else:
                logger.warning(f"Unknown conversion target: {target_format}")
                return []
        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}")
            return []

    def _extract_archive(self, file_path: Path, output_dir: Path) -> List[Path]:
        """Extract archive and return list of extracted files."""
        extracted_files = []
        extract_dir = output_dir / f"extracted_{file_path.stem}"
        extract_dir.mkdir(parents=True, exist_ok=True)

        try:
            if file_path.suffix.lower() == '.zip':
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(extract_dir)
            elif file_path.suffix.lower() in ['.tar', '.gz']:
                with tarfile.open(file_path, 'r:*') as tar_ref:
                    tar_ref.extractall(extract_dir)
            elif file_path.suffix.lower() == '.7z':
                # Requires 7z command line tool
                if shutil.which('7z'):
                    subprocess.run(['7z', 'x', str(file_path), f'-o{extract_dir}'],
                                 check=True, capture_output=True)
                else:
                    logger.warning("7z tool not available for .7z extraction")
                    return []

            # Recursively find all files in extracted directory
            for extracted_file in extract_dir.rglob('*'):
                if extracted_file.is_file():
                    extracted_files.append(extracted_file)

        except Exception as e:
            logger.error(f"Error extracting {file_path}: {e}")

        return extracted_files

    def _convert_to_text(self, file_path: Path, output_dir: Path) -> List[Path]:
        """Convert various formats to text using available tools."""
        output_file = output_dir / f"{file_path.stem}.txt"

        try:
            extension = file_path.suffix.lower()

            if extension in ['.pptx', '.ppt']:
                content = self._extract_powerpoint_text(file_path)
            elif extension == '.rtf':
                content = self._extract_rtf_text(file_path)
            elif extension in ['.odt', '.ods', '.odp']:
                content = self._extract_libreoffice_text(file_path)
            elif extension == '.epub':
                content = self._extract_epub_text(file_path)
            else:
                # Try generic text extraction
                content = self._generic_text_extraction(file_path)

            if content:
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(content)
                return [output_file]

        except Exception as e:
            logger.error(f"Error converting {file_path} to text: {e}")

        return []

    def _extract_powerpoint_text(self, file_path: Path) -> str:
        """Extract text from PowerPoint files."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)

            return "\n".join(text_content)
        except ImportError:
            logger.warning("python-pptx not available for PowerPoint conversion")
            return ""
        except Exception as e:
            logger.error(f"Error extracting PowerPoint text: {e}")
            return ""

    def _extract_rtf_text(self, file_path: Path) -> str:
        """Extract text from RTF files."""
        try:
            from striprtf.striprtf import rtf_to_text
            with open(file_path, 'r', encoding='utf-8') as f:
                rtf_content = f.read()
            return rtf_to_text(rtf_content)
        except ImportError:
            logger.warning("striprtf not available for RTF conversion")
            return ""
        except Exception as e:
            logger.error(f"Error extracting RTF text: {e}")
            return ""

    def _extract_libreoffice_text(self, file_path: Path) -> str:
        """Extract text using LibreOffice command line tools."""
        if not shutil.which('libreoffice'):
            logger.warning("LibreOffice not available for document conversion")
            return ""

        try:
            # Convert to text using LibreOffice headless mode
            result = subprocess.run([
                'libreoffice', '--headless', '--convert-to', 'txt',
                '--outdir', '/tmp', str(file_path)
            ], capture_output=True, text=True, timeout=30)

            if result.returncode == 0:
                txt_file = Path('/tmp') / f"{file_path.stem}.txt"
                if txt_file.exists():
                    content = txt_file.read_text(encoding='utf-8')
                    txt_file.unlink()  # Clean up
                    return content
        except Exception as e:
            logger.error(f"Error using LibreOffice conversion: {e}")

        return ""

    def _extract_epub_text(self, file_path: Path) -> str:
        """Extract text from EPUB files."""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(str(file_path))
            text_content = []

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text_content.append(soup.get_text())

            return "\n".join(text_content)
        except ImportError:
            logger.warning("ebooklib not available for EPUB conversion")
            return ""
        except Exception as e:
            logger.error(f"Error extracting EPUB text: {e}")
            return ""

    def _generic_text_extraction(self, file_path: Path) -> str:
        """Generic text extraction as fallback."""
        try:
            # Try reading as text with different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'ascii']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
        except Exception as e:
            logger.error(f"Generic text extraction failed: {e}")

        return ""


class ContentExtractor:
    """Handles content extraction from various file types."""

    def __init__(self):
        self.supported_extractors = {
            '.txt': self._extract_text,
            '.md': self._extract_text,
            '.py': self._extract_text,
            '.js': self._extract_text,
            '.html': self._extract_text,
            '.css': self._extract_text,
            '.json': self._extract_text,
            '.xml': self._extract_text,
            '.csv': self._extract_csv,
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.doc': self._extract_docx,  # Will try to handle legacy format
            '.xlsx': self._extract_excel,
            '.xls': self._extract_excel,
            '.png': self._extract_image_ocr,
            '.jpg': self._extract_image_ocr,
            '.jpeg': self._extract_image_ocr,
            '.tiff': self._extract_image_ocr,
            '.tif': self._extract_image_ocr,
            '.bmp': self._extract_image_ocr,
            '.gif': self._extract_image_ocr,
        }

    def extract_content(self, file_path: Path) -> Optional[str]:
        """Extract text content from a file."""
        extension = file_path.suffix.lower()
        extractor = self.supported_extractors.get(extension)

        if not extractor:
            logger.warning(f"No extractor available for {extension}: {file_path}")
            return None

        try:
            content = extractor(file_path)
            if content and len(content.strip()) > 0:
                return content.strip()
            else:
                logger.warning(f"No content extracted from: {file_path}")
                return None
        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {e}")
            return None

    def _extract_text(self, file_path: Path) -> str:
        """Extract content from plain text files."""
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise UnicodeDecodeError(f"Could not decode {file_path} with any encoding")

    def _extract_csv(self, file_path: Path) -> str:
        """Extract content from CSV files."""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            logger.warning(f"Failed to parse CSV {file_path}, trying as text: {e}")
            return self._extract_text(file_path)

    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files."""
        text = ""
        try:
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text()
            doc.close()
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {e}")
            raise
        return text

    def _extract_docx(self, file_path: Path) -> str:
        """Extract text from Word documents."""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting DOCX {file_path}: {e}")
            raise

    def _extract_excel(self, file_path: Path) -> str:
        """Extract text from Excel files."""
        try:
            # Read all sheets
            xl_file = pd.ExcelFile(file_path)
            text = ""
            for sheet_name in xl_file.sheet_names:
                df = xl_file.parse(sheet_name)
                text += f"Sheet: {sheet_name}\n"
                text += df.to_string() + "\n\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting Excel {file_path}: {e}")
            raise

    def _extract_image_ocr(self, file_path: Path) -> str:
        """Extract text from images using OCR."""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, config=TESSERACT_CONFIG)
            return text
        except Exception as e:
            logger.error(f"Error performing OCR on {file_path}: {e}")
            raise


class TextChunker:
    """Handles text chunking for vector embeddings."""

    def __init__(self, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP):
        self.chunk_size = chunk_size
        self.overlap = overlap

    def chunk_text(self, text: str, source_file: str) -> List[Dict[str, Any]]:
        """Split text into overlapping chunks."""
        if len(text) <= self.chunk_size:
            return [{
                'text': text,
                'source_file': source_file,
                'chunk_index': 0,
                'total_chunks': 1
            }]

        chunks = []
        start = 0
        chunk_index = 0

        while start < len(text):
            end = start + self.chunk_size

            # Try to break at word boundaries
            if end < len(text):
                # Look for the last space within a reasonable distance
                space_pos = text.rfind(' ', start, end)
                if space_pos > start + self.chunk_size * 0.8:  # Only if we don't lose too much
                    end = space_pos

            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append({
                    'text': chunk_text,
                    'source_file': source_file,
                    'chunk_index': chunk_index,
                    'total_chunks': 0  # Will be updated after all chunks are created
                })
                chunk_index += 1

            start = end - self.overlap if end - self.overlap > start else end

        # Update total_chunks for all chunks
        total_chunks = len(chunks)
        for chunk in chunks:
            chunk['total_chunks'] = total_chunks

        return chunks


class VectorDatabase:
    """Handles FAISS vector database operations."""

    def __init__(self, model_name: str = EMBEDDING_MODEL):
        self.model = SentenceTransformer(model_name)
        self.dimension = self.model.get_sentence_embedding_dimension()
        self.index = faiss.IndexFlatIP(self.dimension)  # Inner product for cosine similarity
        self.chunks = []

    def add_chunks(self, chunks: List[Dict[str, Any]]):
        """Add text chunks to the vector database."""
        if not chunks:
            return

        texts = [chunk['text'] for chunk in chunks]
        logger.info(f"Generating embeddings for {len(texts)} chunks...")

        # Generate embeddings in batches to avoid memory issues
        batch_size = 32
        embeddings = []

        for i in tqdm(range(0, len(texts), batch_size), desc="Generating embeddings"):
            batch_texts = texts[i:i + batch_size]
            batch_embeddings = self.model.encode(batch_texts, normalize_embeddings=True)
            embeddings.extend(batch_embeddings)

        embeddings_array = np.array(embeddings, dtype=np.float32)

        # Add to FAISS index
        self.index.add(embeddings_array)
        self.chunks.extend(chunks)

        logger.info(f"Added {len(chunks)} chunks to vector database. Total chunks: {len(self.chunks)}")

    def save(self, path: Path):
        """Save the vector database to disk."""
        path.mkdir(parents=True, exist_ok=True)

        # Save FAISS index
        faiss.write_index(self.index, str(path / "index.faiss"))

        # Save chunks metadata
        with open(path / "chunks.json", 'w', encoding='utf-8') as f:
            json.dump(self.chunks, f, ensure_ascii=False, indent=2)

        # Save model info
        with open(path / "model_info.json", 'w') as f:
            json.dump({
                'model_name': self.model.get_sentence_embedding_dimension(),
                'dimension': self.dimension,
                'total_chunks': len(self.chunks)
            }, f, indent=2)

        logger.info(f"Vector database saved to {path}")


def load_assessment_report() -> pd.DataFrame:
    """Load the assessment report."""
    if not ASSESSMENT_REPORT.exists():
        raise FileNotFoundError(f"Assessment report not found: {ASSESSMENT_REPORT}")

    return pd.read_csv(ASSESSMENT_REPORT)


def process_files(assessment_df: pd.DataFrame) -> List[Dict[str, Any]]:
    """Process all readable files and extract content, with automatic conversion for incompatible files."""
    extractor = ContentExtractor()
    chunker = TextChunker()
    converter = FileConverter()

    # Create conversion directory
    conversion_dir = PROCESSED_DIR / "converted"
    conversion_dir.mkdir(parents=True, exist_ok=True)

    # Process readable files first
    readable_files = assessment_df[assessment_df['status'] == 'Readable']
    all_chunks = []

    logger.info(f"Processing {len(readable_files)} readable files...")

    for _, row in tqdm(readable_files.iterrows(), total=len(readable_files), desc="Processing readable files"):
        file_path = Path(row['file_path'])

        try:
            content = extractor.extract_content(file_path)
            if content:
                chunks = chunker.chunk_text(content, str(file_path))
                all_chunks.extend(chunks)
                logger.debug(f"Extracted {len(chunks)} chunks from {file_path}")
            else:
                logger.warning(f"No content extracted from {file_path}")

        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            continue

    # Process incompatible files that can be converted
    incompatible_files = assessment_df[assessment_df['status'] == 'Incompatible']
    convertible_files = []

    for _, row in incompatible_files.iterrows():
        file_path = Path(row['file_path'])
        if converter.can_convert(file_path):
            convertible_files.append(file_path)

    if convertible_files:
        logger.info(f"Converting {len(convertible_files)} incompatible files...")

        for file_path in tqdm(convertible_files, desc="Converting files"):
            try:
                converted_files = converter.convert_file(file_path, conversion_dir)

                for converted_file in converted_files:
                    try:
                        content = extractor.extract_content(converted_file)
                        if content:
                            # Use original file path for source attribution but note conversion
                            chunks = chunker.chunk_text(content, f"{file_path} (converted)")
                            all_chunks.extend(chunks)
                            logger.debug(f"Extracted {len(chunks)} chunks from converted {file_path}")
                        else:
                            logger.warning(f"No content extracted from converted {converted_file}")
                    except Exception as e:
                        logger.error(f"Failed to process converted file {converted_file}: {e}")

            except Exception as e:
                logger.error(f"Failed to convert {file_path}: {e}")
                continue

    # Process extracted archive contents
    if conversion_dir.exists():
        for extracted_file in conversion_dir.rglob('*'):
            if extracted_file.is_file() and extracted_file.suffix.lower() in ['.txt']:
                # Skip files we already processed during conversion
                continue

    logger.info(f"Total chunks extracted: {len(all_chunks)}")
    return all_chunks


def main():
    parser = argparse.ArgumentParser(description="Process files and build vector index")
    parser.add_argument("--force", action="store_true", help="Force rebuild even if index exists")
    args = parser.parse_args()

    # Check if processing has already been done
    if VECTOR_DB_PATH.exists() and not args.force:
        print(f"Vector database already exists at {VECTOR_DB_PATH}")
        print("Use --force to rebuild, or proceed to querying with app.py")
        return

    # Load assessment report
    try:
        assessment_df = load_assessment_report()
        logger.info(f"Loaded assessment report with {len(assessment_df)} files")
    except FileNotFoundError as e:
        logger.error(f"Assessment report not found. Please run assess.py first.")
        sys.exit(1)

    # Check if there are readable files
    readable_count = (assessment_df['status'] == 'Readable').sum()
    if readable_count == 0:
        logger.error("No readable files found in assessment report.")
        sys.exit(1)

    logger.info(f"Found {readable_count} readable files to process")

    # Process files
    all_chunks = process_files(assessment_df)

    if not all_chunks:
        logger.error("No content chunks were extracted from any files.")
        sys.exit(1)

    # Build vector database
    logger.info("Building vector database...")
    vector_db = VectorDatabase()
    vector_db.add_chunks(all_chunks)

    # Save vector database
    vector_db.save(VECTOR_DB_PATH)

    logger.info("Processing completed successfully!")
    logger.info(f"Vector database saved to: {VECTOR_DB_PATH}")
    logger.info(f"Total chunks indexed: {len(all_chunks)}")
    logger.info("You can now run the query interface with: python src/app.py")


if __name__ == "__main__":
    main()